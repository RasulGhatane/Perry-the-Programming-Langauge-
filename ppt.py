from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

ppt = Presentation()

def add_slide(title_text, content_text):
    slide_layout = ppt.slide_layouts[5] 
    slide = ppt.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(30, 30, 30)  

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Inches(0.8)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) 

    left, top, width, height = Inches(1), Inches(1.5), Inches(8.5), Inches(5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    text_frame.paragraphs[0].font.size = Inches(0.5)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200) 

    for line in content_text.split("\n"):
        p = text_frame.add_paragraph()
        p.text = "• " + line.strip()
        p.font.size = Inches(0.4)
        p.font.color.rgb = RGBColor(220, 220, 220)

add_slide("Understanding rizz.c", "A breakdown of the custom tokenization and parsing system")

add_slide("Overview",
          "A custom tokenizer and parser written in C\n"
          "Implements a simple programming language\n"
          "Handles input, variable declaration, loops, and output")

add_slide("Features",
          "Tokenization of a custom language\n"
          "Parsing and validation of tokens\n"
          "Execution of basic programming constructs\n"
          "Error handling and memory management")

add_slide("Data Structures",
          "Token - Represents different tokens in the language\n"
          "Value - Stores different types of values\n"
          "ProgramState - Maintains global execution state")

add_slide("Lexical Analysis",
          "Tokenizer scans source code and generates tokens\n"
          "Recognizes integers, floats, strings, and identifiers\n"
          "Supports keywords: yap (print), cook (input), sigma (declare), gyatt (loop)")

add_slide("Parsing & Execution",
          "Checks token sequences for validity\n"
          "Executes commands like print, input, and loops\n"
          "Stores and retrieves variable values")

add_slide("Error Handling & Memory Management",
          "Uses safe_malloc to prevent memory allocation failures\n"
          "Implements runtime_error function for error reporting\n"
          "Properly cleans up allocated memory after execution")

add_slide("Main Function Workflow",
          "Reads source code from a file\n"
          "Tokenizes the input\n"
          "Parses and validates tokens\n"
          "Executes the corresponding actions\n"
          "Cleans up memory before exiting")

ppt_path = r"C:\Users\rasul\OneDrive\Documents\rizz_presentation.pptx"
ppt.save(ppt_path)

ppt_path
